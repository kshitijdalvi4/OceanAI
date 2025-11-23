from fastapi import FastAPI, Depends, HTTPException, status, File, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from pydantic import BaseModel, EmailStr
from typing import List, Optional, Dict, Any
from datetime import datetime, timedelta
from fastapi import Body
import os
import json
import re
import time
import io
import uuid

# Google Gemini
from google import genai
from google.genai import types

# LangChain for RAG
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_core.embeddings import Embeddings
import chromadb

# Database
from sqlalchemy import create_engine, Column, Integer, String, DateTime, Text, JSON, ForeignKey, Float, Boolean
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, Session, relationship

# Document Generation
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt

# Auth
import jwt
import bcrypt

# ============================================================================
# CONFIGURATION
# ============================================================================

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "AIzaSyCTmZaq2CM-5EOS3YYrhFosl5yweDmo0F4")
SECRET_KEY = os.getenv("SECRET_KEY", "your-secret-key-change-in-production-oceanai-2025")
DATABASE_URL = os.getenv("DATABASE_URL", "sqlite:///./oceanai.db")
CHROMA_PATH = os.getenv("CHROMA_PATH", "./chroma_db")

# Initialize Gemini Client
gemini_client = genai.Client(api_key=GEMINI_API_KEY)

# ============================================================================
# FASTAPI APP SETUP
# ============================================================================

app = FastAPI(
    title="OceanAI Document Generation Platform",
    description="AI-powered document authoring with RAG and ContentMemoryBuffer",
    version="1.0.0"
)

# CORS Middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:8000",
        "http://localhost:5173", 
        "http://localhost:5174",
        "https://codemos-services.co.in",
        "http://codemos-services.co.in"
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ============================================================================
# DATABASE MODELS (SQLAlchemy)
# ============================================================================

Base = declarative_base()
engine = create_engine(DATABASE_URL, connect_args={"check_same_thread": False} if "sqlite" in DATABASE_URL else {})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)

class User(Base):
    __tablename__ = "users"
    
    id = Column(Integer, primary_key=True, index=True)
    email = Column(String, unique=True, index=True, nullable=False)
    hashed_password = Column(String, nullable=False)
    full_name = Column(String, nullable=True)
    created_at = Column(DateTime, default=datetime.utcnow)
    last_login = Column(DateTime, nullable=True)
    
    # Relationships
    projects = relationship("Project", back_populates="owner", cascade="all, delete-orphan")

class Project(Base):
    __tablename__ = "projects"
    
    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"), nullable=False)
    title = Column(String, nullable=False)
    document_type = Column(String, nullable=False)  # 'docx' or 'pptx'
    main_topic = Column(Text, nullable=False)
    configuration = Column(JSON, default={})  # Stores outline/slides structure
    status = Column(String, default="draft")  # draft, generated, completed
    created_at = Column(DateTime, default=datetime.utcnow)
    updated_at = Column(DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)
    
    # Relationships
    owner = relationship("User", back_populates="projects")
    contents = relationship("Content", back_populates="project", cascade="all, delete-orphan")
    refinements = relationship("Refinement", back_populates="project", cascade="all, delete-orphan")

class Content(Base):
    __tablename__ = "contents"
    
    id = Column(Integer, primary_key=True, index=True)
    project_id = Column(Integer, ForeignKey("projects.id"), nullable=False)
    section_id = Column(String, nullable=False)  # section-1, slide-1, etc.
    title = Column(String, nullable=False)
    content_text = Column(Text, nullable=False)
    version = Column(Integer, default=1)
    is_current = Column(Boolean, default=True)
    generated_at = Column(DateTime, default=datetime.utcnow)
    
    # Relationships
    project = relationship("Project", back_populates="contents")

class Refinement(Base):
    __tablename__ = "refinements"
    
    id = Column(Integer, primary_key=True, index=True)
    project_id = Column(Integer, ForeignKey("projects.id"), nullable=False)
    section_id = Column(String, nullable=False)
    refinement_prompt = Column(Text, nullable=True)
    feedback_type = Column(String, nullable=True)  # 'like', 'dislike', 'neutral'
    comment = Column(Text, nullable=True)
    old_content = Column(Text, nullable=True)
    new_content = Column(Text, nullable=True)
    created_at = Column(DateTime, default=datetime.utcnow)
    
    # Relationships
    project = relationship("Project", back_populates="refinements")

# Create all tables
Base.metadata.create_all(bind=engine)

# ============================================================================
# CUSTOM GEMINI EMBEDDINGS CLASS
# ============================================================================

class GeminiEmbeddings(Embeddings):
    """Custom embeddings class using Gemini API"""
    
    def __init__(self, client):
        self.client = client
    
    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        """Embed a list of documents"""
        embeddings = []
        for text in texts:
            try:
                response = self.client.models.embed_content(
                    model="models/text-embedding-004",
                    content=text
                )
                embeddings.append(response.embedding)
                time.sleep(0.1)  # Rate limiting
            except Exception as e:
                print(f"Error embedding document: {e}")
                embeddings.append([0.0] * 768)
        return embeddings
    
    def embed_query(self, text: str) -> List[float]:
        """Embed a single query"""
        try:
            response = self.client.models.embed_content(
                model="models/text-embedding-004",
                content=text
            )
            return response.embedding
        except Exception as e:
            print(f"Error embedding query: {e}")
            return [0.0] * 768

# Initialize embeddings
embeddings = GeminiEmbeddings(gemini_client)

# ============================================================================
# CHROMADB SETUP FOR RAG
# ============================================================================

os.makedirs(CHROMA_PATH, exist_ok=True)
chroma_client = chromadb.PersistentClient(path=CHROMA_PATH)

# ============================================================================
# PYDANTIC MODELS (API Schemas)
# ============================================================================

class UserRegister(BaseModel):
    email: EmailStr
    password: str
    full_name: Optional[str] = None

class UserLogin(BaseModel):
    email: EmailStr
    password: str

class Token(BaseModel):
    access_token: str
    token_type: str
    user_email: str
    user_id: int

class ProjectCreate(BaseModel):
    title: str
    document_type: str  # 'docx' or 'pptx'
    main_topic: str

class ProjectUpdate(BaseModel):
    title: Optional[str] = None
    main_topic: Optional[str] = None
    configuration: Optional[Dict[str, Any]] = None

class ConfigurationUpdate(BaseModel):
    configuration: Dict[str, Any]

class GenerateRequest(BaseModel):
    section_id: Optional[str] = None  # If None, generate all

class RefineRequest(BaseModel):
    section_id: str
    prompt: str

class FeedbackRequest(BaseModel):
    section_id: str
    feedback_type: str  # 'like' or 'dislike'
    comment: Optional[str] = None

class ChatQuery(BaseModel):
    project_id: int
    question: str

# NEW: Pydantic model for outline suggestion
class OutlineSuggestionRequest(BaseModel):
    topic: str
    document_type: str

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

security = HTTPBearer()

def get_db():
    """Database dependency"""
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

def hash_password(password: str) -> str:
    """Hash password using bcrypt"""
    return bcrypt.hashpw(password.encode(), bcrypt.gensalt()).decode()

def verify_password(password: str, hashed: str) -> bool:
    """Verify password against hash"""
    return bcrypt.checkpw(password.encode(), hashed.encode())

def create_access_token(user_id: int, email: str) -> str:
    """Create JWT token"""
    payload = {
        "user_id": user_id,
        "email": email,
        "exp": datetime.utcnow() + timedelta(days=7)
    }
    return jwt.encode(payload, SECRET_KEY, algorithm="HS256")

def verify_token(credentials: HTTPAuthorizationCredentials = Depends(security)) -> dict:
    """Verify JWT token and return user info"""
    try:
        payload = jwt.decode(credentials.credentials, SECRET_KEY, algorithms=["HS256"])
        return payload
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token has expired")
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail="Invalid token")

def retry_on_503(max_retries=3, delay=2):
    """Decorator to retry on 503 errors"""
    def decorator(func):
        def wrapper(*args, **kwargs):
            for attempt in range(max_retries):
                try:
                    return func(*args, **kwargs)
                except Exception as e:
                    error_str = str(e)
                    if '503' in error_str and attempt < max_retries - 1:
                        print(f"503 error, retrying in {delay}s... (attempt {attempt + 1}/{max_retries})")
                        time.sleep(delay * (attempt + 1))
                        continue
                    raise
            return func(*args, **kwargs)
        return wrapper
    return decorator

def clean_json_response(text: str) -> str:
    """Clean Gemini response to extract JSON"""
    # Remove markdown code blocks
    text = re.sub(r'```json\s*', '', text)
    text = re.sub(r'```\s*', '', text)
    text = text.strip()
    
    # Try to find JSON array first (for outline suggestions)
    array_match = re.search(r'\[\s*\{.*?\}\s*\]', text, re.DOTALL)
    if array_match:
        json_str = array_match.group(0)
        # Clean trailing commas
        json_str = re.sub(r',\s*}', '}', json_str)
        json_str = re.sub(r',\s*]', ']', json_str)
        return json_str
    
    # Try to find JSON object
    json_match = re.search(r'\{[^{}]*(?:\{[^{}]*\}[^{}]*)*\}', text, re.DOTALL)
    if json_match:
        json_str = json_match.group(0)
        json_str = re.sub(r',\s*}', '}', json_str)
        json_str = re.sub(r',\s*]', ']', json_str)
        return json_str
    
    # If starts with [ or {, try to find the end
    if text.startswith('['):
        end_idx = text.rfind(']')
        if end_idx != -1:
            text = text[:end_idx + 1]
    elif text.startswith('{'):
        end_idx = text.rfind('}')
        if end_idx != -1:
            text = text[:end_idx + 1]
    
    return text.strip()

@retry_on_503(max_retries=3, delay=2)
def call_gemini(prompt: str, model: str = "gemini-2.0-flash-lite") -> str:
    """Call Gemini API with retry logic"""
    response = gemini_client.models.generate_content(
        model=model,
        contents=prompt,
    )
    return response.text

# ============================================================================
# CONTENT MEMORY BUFFER (RAG System)
# ============================================================================

class ContentMemoryBuffer:
    """
    Manages context and memory for document generation using RAG.
    Stores all generated content in ChromaDB for context-aware refinements.
    """
    
    def __init__(self, project_id: int):
        self.project_id = project_id
        self.collection_name = f"project_{project_id}"
        
        try:
            self.collection = chroma_client.get_collection(self.collection_name)
        except:
            self.collection = chroma_client.create_collection(
                name=self.collection_name,
                metadata={"project_id": str(project_id)}
            )
    
    def add_content(self, section_id: str, title: str, content: str, metadata: Dict):
        """Add generated content to memory buffer"""
        try:
            doc_id = f"{self.project_id}_{section_id}_{uuid.uuid4().hex[:8]}"
            
            # Get embedding
            embedding = embeddings.embed_query(content)
            
            # Add to collection
            self.collection.add(
                documents=[content],
                embeddings=[embedding],
                metadatas=[{
                    "section_id": section_id,
                    "title": title,
                    "timestamp": datetime.utcnow().isoformat(),
                    **metadata
                }],
                ids=[doc_id]
            )
            print(f"✓ Added content to memory: {section_id}")
        except Exception as e:
            print(f"Error adding content to memory: {e}")
    
    def query_context(self, query: str, n_results: int = 3) -> List[Dict]:
        """Query relevant context from memory buffer"""
        try:
            results = self.collection.query(
                query_texts=[query],
                n_results=n_results
            )
            
            if results and results['documents']:
                contexts = []
                for i in range(len(results['documents'][0])):
                    contexts.append({
                        'content': results['documents'][0][i],
                        'metadata': results['metadatas'][0][i] if results['metadatas'] else {}
                    })
                return contexts
            return []
        except Exception as e:
            print(f"Error querying context: {e}")
            return []
    
    def get_project_context(self) -> str:
        """Get full project context"""
        try:
            all_docs = self.collection.get()
            if not all_docs['documents']:
                return ""
            # Return last 10 documents
            return "\n\n".join(all_docs['documents'][-10:])
        except Exception as e:
            print(f"Error getting project context: {e}")
            return ""

# ============================================================================
# RAG SERVICE
# ============================================================================

class RAGService:
    """Service for RAG-based content generation and refinement"""
    
    @staticmethod
    def generate_section_content(
        title: str,
        topic: str,
        section_type: str,
        context: str = ""
    ) -> str:
        """Generate content for a section using RAG"""
        
        system_prompt = """You are a professional business document writer specializing in creating high-quality, well-structured content."""
        
        context_part = f"\n\nPrevious Context:\n{context}" if context else ""
        
        prompt = f"""{system_prompt}

Generate professional content for the following section.

DOCUMENT TOPIC: {topic}
SECTION TITLE: {title}
SECTION TYPE: {section_type}
{context_part}

Requirements:
- Write clear, professional content (200-400 words)
- Use proper paragraphs and structure
- Be specific and informative
- Maintain consistency with previous sections
- Write in plain text without markdown formatting

Content:"""
        
        return call_gemini(prompt)
    
    @staticmethod
    def refine_content(
        original: str,
        refinement_prompt: str,
        context: List[Dict]
    ) -> str:
        """Refine content based on user prompt with RAG context"""
        
        context_text = "\n\n".join([f"Related section: {c['content'][:200]}..." for c in context])
        
        prompt = f"""You are refining a section of a business document.

Previous related content:
{context_text}

Current content to refine:
{original}

User's refinement request:
{refinement_prompt}

Provide the refined content below. Write in plain text without markdown formatting.

Refined content:"""
        
        return call_gemini(prompt)
    
    @staticmethod
    def suggest_outline(topic: str, document_type: str) -> List[Dict]:
        """Generate document outline using AI"""
        
        section_type = "sections" if document_type == "docx" else "slides"
        count = 6 if document_type == "docx" else 8
        
        prompt = f"""You are a professional document outline generator.

Create a detailed outline for a {document_type.upper()} document about: {topic}

You MUST create exactly {count} {section_type}. Each section must have a title and description.

Return ONLY a valid JSON array with exactly {count} items, no markdown, no explanation:
[
  {{"title": "Section Title 1", "description": "What this section covers"}},
  {{"title": "Section Title 2", "description": "What this section covers"}},
  ...
]

Generate exactly {count} {section_type} now:"""
        
        response = call_gemini(prompt)
        
        print(f"🤖 Raw Gemini response: {response[:500]}")  # Debug
        
        try:
            # Try to extract JSON array from response
            cleaned = clean_json_response(response)
            print(f"🧹 Cleaned response: {cleaned[:500]}")  # Debug
            
            suggestions = json.loads(cleaned)
            
            # Validate it's a list
            if not isinstance(suggestions, list):
                raise ValueError("Response is not a list")
            
            # Ensure we have enough suggestions
            if len(suggestions) < count:
                print(f"⚠️  Only got {len(suggestions)} suggestions, expected {count}")
                # Generate additional ones
                for i in range(len(suggestions), count):
                    suggestions.append({
                        "title": f"Section {i+1}",
                        "description": f"Additional content for {topic}"
                    })
            
            print(f"✅ Successfully parsed {len(suggestions)} suggestions")
            return suggestions[:count]  # Return exactly the count we need
            
        except Exception as e:
            print(f"❌ Error parsing outline: {e}")
            print(f"📄 Full response was: {response}")
            
            # Fallback: Generate default outline
            print(f"🔄 Using fallback outline with {count} sections")
            return [
                {"title": f"Section {i+1}: {topic}", "description": f"Content about {topic} - part {i+1}"}
                for i in range(count)
            ]

# ============================================================================
# API ROUTES - AUTHENTICATION
# ============================================================================

@app.get("/")
async def root():
    return {
        "message": "OceanAI Document Generation Platform",
        "status": "running",
        "version": "1.0.0",
        "model": "Gemini 2.0 Flash Lite",
        "embeddings": "Gemini text-embedding-004",
        "features": ["RAG", "ContentMemoryBuffer", "Document Export", "AI Outline Suggest Direct"]
    }

@app.post("/api/auth/register", response_model=Token)
async def register(user_data: UserRegister, db: Session = Depends(get_db)):
    """Register a new user"""
    
    # Check if user exists
    existing = db.query(User).filter(User.email == user_data.email).first()
    if existing:
        raise HTTPException(status_code=400, detail="Email already registered")
    
    # Create new user
    hashed_pw = hash_password(user_data.password)
    new_user = User(
        email=user_data.email,
        hashed_password=hashed_pw,
        full_name=user_data.full_name
    )
    
    db.add(new_user)
    db.commit()
    db.refresh(new_user)
    
    # Create token
    token = create_access_token(new_user.id, new_user.email)
    
    return Token(
        access_token=token,
        token_type="bearer",
        user_email=new_user.email,
        user_id=new_user.id
    )

@app.post("/api/auth/login", response_model=Token)
async def login(credentials: UserLogin, db: Session = Depends(get_db)):
    """Login user"""
    
    user = db.query(User).filter(User.email == credentials.email).first()
    
    if not user or not verify_password(credentials.password, user.hashed_password):
        raise HTTPException(status_code=401, detail="Invalid email or password")
    
    # Update last login
    user.last_login = datetime.utcnow()
    db.commit()
    
    # Create token
    token = create_access_token(user.id, user.email)
    
    return Token(
        access_token=token,
        token_type="bearer",
        user_email=user.email,
        user_id=user.id
    )

# ============================================================================
# API ROUTES - PROJECT MANAGEMENT
# ============================================================================

@app.get("/api/projects")
async def list_projects(
    token_data: dict = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Get all projects for authenticated user"""
    
    user_id = token_data["user_id"]
    projects = db.query(Project).filter(Project.user_id == user_id).order_by(Project.updated_at.desc()).all()
    
    return {"projects": projects, "total": len(projects)}

@app.post("/api/projects")
async def create_project(
    project_data: ProjectCreate,
    token_data: dict = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Create a new project"""
    
    user_id = token_data["user_id"]
    
    new_project = Project(
        user_id=user_id,
        title=project_data.title,
        document_type=project_data.document_type,
        main_topic=project_data.main_topic,
        configuration={}
    )
    
    db.add(new_project)
    db.commit()
    db.refresh(new_project)
    
    return {
        "message": "Project created successfully",
        "project": new_project
    }

@app.get("/api/projects/{project_id}")
async def get_project(
    project_id: int,
    token_data: dict = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Get project details with contents"""
    
    user_id = token_data["user_id"]
    
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == user_id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Get current contents
    contents = db.query(Content).filter(
        Content.project_id == project_id,
        Content.is_current == True
    ).order_by(Content.section_id).all()
    
    return {
        "project": project,
        "contents": contents
    }

@app.put("/api/projects/{project_id}")
async def update_project(
    project_id: int,
    update_data: ProjectUpdate,
    token_data: dict = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Update project details"""
    
    user_id = token_data["user_id"]
    
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == user_id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    if update_data.title:
        project.title = update_data.title
    if update_data.main_topic:
        project.main_topic = update_data.main_topic
    if update_data.configuration is not None:
        project.configuration = update_data.configuration
    
    project.updated_at = datetime.utcnow()
    db.commit()
    
    return {"message": "Project updated successfully", "project": project}

@app.put("/api/projects/{project_id}/config")
async def update_config(
    project_id: int,
    config: ConfigurationUpdate,
    token_data: dict = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Update project configuration (outline/slides)"""
    
    user_id = token_data["user_id"]
    
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == user_id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    project.configuration = config.configuration
    project.updated_at = datetime.utcnow()
    db.commit()
    
    return {"message": "Configuration updated successfully"}

@app.delete("/api/projects/{project_id}")
async def delete_project(
    project_id: int,
    token_data: dict = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Delete a project"""
    
    user_id = token_data["user_id"]
    
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == user_id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Delete from ChromaDB
    try:
        chroma_client.delete_collection(f"project_{project_id}")
    except:
        pass
    
    db.delete(project)
    db.commit()
    
    return {"message": "Project deleted successfully"}

# ============================================================================
# API ROUTES - CONTENT GENERATION
# ============================================================================

@app.post("/api/projects/{project_id}/generate")
async def generate_content(
    project_id: int,
    request: GenerateRequest,
    token_data: dict = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Generate AI content for project sections"""
    
    user_id = token_data["user_id"]
    
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == user_id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Initialize memory buffer
    memory_buffer = ContentMemoryBuffer(project_id)
    
    # Get configuration
    config = project.configuration
    
    # Auto-create BASIC default sections if none exist (REDUCED TO 3)
    if not config or (not config.get('sections') and not config.get('slides')):
        if project.document_type == 'docx':
            config = {
                'sections': [
                    {'title': 'Introduction', 'description': 'Opening remarks'},
                    {'title': 'Body', 'description': 'Main content'},
                    {'title': 'Conclusion', 'description': 'Closing thoughts'}
                ]
            }
        else:  # pptx
            config = {
                'slides': [
                    {'title': 'Title Slide', 'description': 'Cover page'},
                    {'title': 'Main Points', 'description': 'Key information'},
                    {'title': 'Summary', 'description': 'Wrap up'}
                ]
            }
        
        # Update project config
        project.configuration = config
        db.commit()
        
        print(f"⚠️  Using BASIC DEFAULT outline (3 sections) for project {project_id}")
    
    sections = config.get('sections', []) if project.document_type == 'docx' else config.get('slides', [])
    
    if not sections:
        raise HTTPException(status_code=400, detail="No sections/slides configured")
    
    generated = []
    
    for i, section in enumerate(sections):
        section_id = f"{'section' if project.document_type == 'docx' else 'slide'}-{i+1}"
        
        # Skip if specific section requested
        if request.section_id and section_id != request.section_id:
            continue
        
        title = section.get('title', f"Section {i+1}")
        
        # Get context from previous sections
        context = memory_buffer.get_project_context()
        
        # Generate content
        content_text = RAGService.generate_section_content(
            title=title,
            topic=project.main_topic,
            section_type=project.document_type,
            context=context
        )
        
        # Mark old content as not current
        db.query(Content).filter(
            Content.project_id == project_id,
            Content.section_id == section_id
        ).update({"is_current": False})
        
        # Save new content
        new_content = Content(
            project_id=project_id,
            section_id=section_id,
            title=title,
            content_text=content_text,
            version=1,
            is_current=True
        )
        db.add(new_content)
        
        # Add to memory buffer
        memory_buffer.add_content(
            section_id=section_id,
            title=title,
            content=content_text,
            metadata={"type": "generation", "version": 1}
        )
        
        generated.append({
            "section_id": section_id,
            "title": title,
            "content": content_text
        })
        
        if request.section_id:
            break
    
    # Update project status
    project.status = "generated"
    project.updated_at = datetime.utcnow()
    db.commit()
    
    return {
        "message": "Content generated successfully",
        "generated": generated,
        "total": len(generated)
    }

@app.post("/api/projects/{project_id}/refine")
async def refine_content(
    project_id: int,
    request: RefineRequest,
    token_data: dict = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Refine specific section content using RAG"""
    
    user_id = token_data["user_id"]
    
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == user_id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Get current content
    current_content = db.query(Content).filter(
        Content.project_id == project_id,
        Content.section_id == request.section_id,
        Content.is_current == True
    ).first()
    
    if not current_content:
        raise HTTPException(status_code=404, detail="Section content not found")
    
    # Get context from memory buffer
    memory_buffer = ContentMemoryBuffer(project_id)
    context_results = memory_buffer.query_context(request.prompt, n_results=3)
    
    # Refine content
    refined_text = RAGService.refine_content(
        original=current_content.content_text,
        refinement_prompt=request.prompt,
        context=context_results
    )
    
    # Mark old content as not current
    current_content.is_current = False
    
    # Create new version
    new_version = current_content.version + 1
    new_content = Content(
        project_id=project_id,
        section_id=request.section_id,
        title=current_content.title,
        content_text=refined_text,
        version=new_version,
        is_current=True
    )
    db.add(new_content)
    
    # Log refinement
    refinement = Refinement(
        project_id=project_id,
        section_id=request.section_id,
        refinement_prompt=request.prompt,
        old_content=current_content.content_text,
        new_content=refined_text
    )
    db.add(refinement)
    
    # Update memory buffer
    memory_buffer.add_content(
        section_id=request.section_id,
        title=current_content.title,
        content=refined_text,
        metadata={"type": "refinement", "version": new_version, "prompt": request.prompt}
    )
    
    project.updated_at = datetime.utcnow()
    db.commit()
    
    return {
        "message": "Content refined successfully",
        "section_id": request.section_id,
        "content": refined_text,
        "version": new_version
    }

@app.post("/api/projects/{project_id}/feedback")
async def submit_feedback(
    project_id: int,
    request: FeedbackRequest,
    token_data: dict = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Submit feedback (like/dislike) for a section"""
    
    user_id = token_data["user_id"]
    
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == user_id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Create feedback record
    feedback = Refinement(
        project_id=project_id,
        section_id=request.section_id,
        feedback_type=request.feedback_type,
        comment=request.comment
    )
    db.add(feedback)
    db.commit()
    
    return {
        "message": "Feedback submitted successfully",
        "feedback_type": request.feedback_type
    }

@app.get("/api/projects/{project_id}/history/{section_id}")
async def get_section_history(
    project_id: int,
    section_id: str,
    token_data: dict = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Get version history for a section"""
    
    user_id = token_data["user_id"]
    
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == user_id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Get all versions
    versions = db.query(Content).filter(
        Content.project_id == project_id,
        Content.section_id == section_id
    ).order_by(Content.version.desc()).all()
    
    # Get refinements
    refinements = db.query(Refinement).filter(
        Refinement.project_id == project_id,
        Refinement.section_id == section_id
    ).order_by(Refinement.created_at.desc()).all()
    
    return {
        "section_id": section_id,
        "versions": versions,
        "refinements": refinements
    }

# ============================================================================
# API ROUTES - AI FEATURES (FIXED DIRECT ENDPOINT)
# ============================================================================

@app.post("/api/suggest-outline-direct")
async def suggest_outline_direct(
    request: OutlineSuggestionRequest,
    token_data: dict = Depends(verify_token)
):
    """
    AI-generate outline WITHOUT creating a project (STANDALONE)
    This is the FIXED endpoint with proper Pydantic model.
    """
    
    print(f"🎯 AI Suggest Outline Direct called for: {request.topic} ({request.document_type})")
    
    try:
        # Generate outline using RAG Service
        suggestions = RAGService.suggest_outline(
            topic=request.topic,
            document_type=request.document_type
        )
        
        print(f"✨ Generated {len(suggestions)} suggestions")
        
        return {
            "success": True,
            "message": "Outline generated successfully",
            "suggestions": suggestions,
            "document_type": request.document_type,
            "count": len(suggestions)
        }
    except Exception as e:
        print(f"❌ Error generating outline: {e}")
        raise HTTPException(
            status_code=500,
            detail=f"Failed to generate outline: {str(e)}"
        )

@app.post("/api/projects/{project_id}/suggest-outline")
async def suggest_outline(
    project_id: int,
    token_data: dict = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """AI-generate document outline (LEGACY - kept for backwards compatibility)"""
    
    user_id = token_data["user_id"]
    
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == user_id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Generate outline
    suggestions = RAGService.suggest_outline(
        topic=project.main_topic,
        document_type=project.document_type
    )
    
    return {
        "message": "Outline generated successfully",
        "suggestions": suggestions,
        "document_type": project.document_type
    }

@app.post("/api/projects/{project_id}/chat")
async def chat_with_project(
    project_id: int,
    query: ChatQuery,
    token_data: dict = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Chat with project content using RAG"""
    
    user_id = token_data["user_id"]
    
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == user_id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Get context from memory buffer
    memory_buffer = ContentMemoryBuffer(project_id)
    context_results = memory_buffer.query_context(query.question, n_results=5)
    
    if not context_results:
        return {
            "question": query.question,
            "answer": "No content available yet. Please generate content first."
        }
    
    # Build context
    context_text = "\n\n".join([
        f"Section: {c['metadata'].get('title', 'Unknown')}\n{c['content']}"
        for c in context_results
    ])
    
    prompt = f"""You are analyzing a business document project.

Project Topic: {project.main_topic}
Document Type: {project.document_type.upper()}

Relevant sections:
{context_text}

User Question: {query.question}

Provide a helpful answer based on the document content. Write in plain text without markdown formatting.

Answer:"""
    
    answer = call_gemini(prompt)
    
    return {
        "question": query.question,
        "answer": answer,
        "sources": [c['metadata'].get('title', 'Unknown') for c in context_results]
    }

# ============================================================================
# API ROUTES - DOCUMENT EXPORT
# ============================================================================

@app.get("/api/projects/{project_id}/export")
async def export_document(
    project_id: int,
    token_data: dict = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Export project as .docx or .pptx file"""
    
    user_id = token_data["user_id"]
    
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == user_id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Get current contents
    contents = db.query(Content).filter(
        Content.project_id == project_id,
        Content.is_current == True
    ).order_by(Content.section_id).all()
    
    if not contents:
        raise HTTPException(status_code=400, detail="No content to export")
    
    if project.document_type == 'docx':
        # Create Word document
        doc = Document()
        
        # Title page
        title = doc.add_heading(project.title, 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        doc.add_paragraph(f"Topic: {project.main_topic}")
        doc.add_paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y')}")
        doc.add_page_break()
        
        # Add sections
        for content in contents:
            doc.add_heading(content.title, level=1)
            
            # Add content paragraphs
            paragraphs = content.content_text.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    p = doc.add_paragraph(para.strip())
                    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            
            doc.add_paragraph()  # Spacing
        
        # Save to buffer
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        
        filename = f"{project.title.replace(' ', '_')}.docx"
        media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        
    else:  # pptx
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = project.title
        title_slide.placeholders[1].text = f"{project.main_topic}\n{datetime.now().strftime('%B %d, %Y')}"
        
        # Content slides
        for content in contents:
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            
            # Title
            slide.shapes.title.text = content.title
            
            # Content
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            # Add content as paragraphs
            paragraphs = content.content_text.split('\n\n')
            for i, para in enumerate(paragraphs):
                if para.strip():
                    if i == 0:
                        text_frame.text = para.strip()
                    else:
                        p = text_frame.add_paragraph()
                        p.text = para.strip()
                        p.level = 0
        
        # Save to buffer
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        filename = f"{project.title.replace(' ', '_')}.pptx"
        media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    
    return StreamingResponse(
        buffer,
        media_type=media_type,
        headers={"Content-Disposition": f"attachment; filename={filename}"}
    )

# ============================================================================
# API ROUTES - ANALYTICS & INSIGHTS
# ============================================================================

@app.get("/api/projects/{project_id}/analytics")
async def get_project_analytics(
    project_id: int,
    token_data: dict = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Get project analytics and statistics"""
    
    user_id = token_data["user_id"]
    
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == user_id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    # Count sections
    sections_count = len(project.configuration.get('sections', []) if project.document_type == 'docx' else project.configuration.get('slides', []))
    
    # Count generated content
    content_count = db.query(Content).filter(
        Content.project_id == project_id,
        Content.is_current == True
    ).count()
    
    # Count refinements
    refinement_count = db.query(Refinement).filter(
        Refinement.project_id == project_id
    ).count()
    
    # Count feedback
    feedback_likes = db.query(Refinement).filter(
        Refinement.project_id == project_id,
        Refinement.feedback_type == 'like'
    ).count()
    
    feedback_dislikes = db.query(Refinement).filter(
        Refinement.project_id == project_id,
        Refinement.feedback_type == 'dislike'
    ).count()
    
    # Total word count
    contents = db.query(Content).filter(
        Content.project_id == project_id,
        Content.is_current == True
    ).all()
    
    total_words = sum(len(c.content_text.split()) for c in contents)
    
    return {
        "project_id": project_id,
        "sections_configured": sections_count,
        "sections_generated": content_count,
        "total_refinements": refinement_count,
        "feedback_likes": feedback_likes,
        "feedback_dislikes": feedback_dislikes,
        "total_words": total_words,
        "completion_percentage": round((content_count / sections_count * 100) if sections_count > 0 else 0, 2)
    }

@app.get("/api/user/stats")
async def get_user_stats(
    token_data: dict = Depends(verify_token),
    db: Session = Depends(get_db)
):
    """Get user statistics"""
    
    user_id = token_data["user_id"]
    
    total_projects = db.query(Project).filter(Project.user_id == user_id).count()
    
    total_documents = db.query(Content).join(Project).filter(
        Project.user_id == user_id,
        Content.is_current == True
    ).count()
    
    total_refinements = db.query(Refinement).join(Project).filter(
        Project.user_id == user_id
    ).count()
    
    # Recent projects
    recent_projects = db.query(Project).filter(
        Project.user_id == user_id
    ).order_by(Project.updated_at.desc()).limit(5).all()
    
    return {
        "total_projects": total_projects,
        "total_documents": total_documents,
        "total_refinements": total_refinements,
        "recent_projects": recent_projects
    }

# ============================================================================
# HEALTH CHECK & SYSTEM
# ============================================================================

@app.get("/health")
async def health_check():
    """System health check"""
    return {
        "status": "healthy",
        "model": "Gemini 2.0 Flash Lite",
        "embeddings": "Gemini text-embedding-004",
        "database": "Connected",
        "chroma_path": CHROMA_PATH,
        "timestamp": datetime.utcnow().isoformat()
    }

@app.get("/api/system/info")
async def system_info():
    """Get system information"""
    return {
        "platform": "OceanAI Document Generation",
        "version": "1.0.0",
        "features": [
            "User Authentication (JWT)",
            "Project Management",
            "RAG-based Content Generation",
            "ContentMemoryBuffer",
            "Iterative Refinement",
            "Feedback System",
            "Document Export (.docx, .pptx)",
            "AI Outline Suggestion (Direct - Fixed)",
            "Project Analytics",
            "Chat with Documents"
        ],
        "supported_formats": ["docx", "pptx"],
        "max_sections": 50,
        "ai_outline": {
            "endpoint": "/api/suggest-outline-direct",
            "method": "POST",
            "body": {"topic": "string", "document_type": "docx or pptx"},
            "default_sections_docx": 3,
            "default_sections_pptx": 3,
            "ai_sections_docx": "6-8",
            "ai_sections_pptx": "8-10"
        }
    }

# ============================================================================
# RUN SERVER
# ============================================================================

if __name__ == "__main__":
    import uvicorn
    
    print("=" * 70)
    print("🌊 OceanAI Document Generation Platform")
    print("=" * 70)
    print("✓ FastAPI backend initialized")
    print("✓ Gemini 2.0 Flash Lite connected")
    print("✓ RAG pipeline with ChromaDB ready")
    print("✓ ContentMemoryBuffer active")
    print("✓ Database: SQLAlchemy ORM")
    print("✓ Authentication: JWT with bcrypt")
    print("✓ FIXED: Direct AI Outline Suggestion with Pydantic model")
    print("✓ Default sections reduced to 3 (vs AI 6-8)")
    print("=" * 70)
    print("🚀 Starting server on http://0.0.0.0:8000")
    print("📚 API docs at http://0.0.0.0:8000/docs")
    print("=" * 70)
    
    uvicorn.run(app, host="0.0.0.0", port=8000)